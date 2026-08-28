from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Q3 2026 Distillation Unit Inspection Report"
subtitle.text = "Detailed Findings"

# Summary Slide
summary_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(summary_slide_layout)
title = slide.shapes.title
body_shape = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(6.0), Inches(5.0))

title.text = "Summary of Detailed Findings"

# Add bullet points to the body shape
bullet_points = [
    "Heat Exchangers (HX-101 to HX-104): No significant fouling detected. Thickness readings are within the 2mm corrosion allowance.",
    "Piping Network (Line L-205): Minor surface oxidation observed on the external insulation cladding. Recommended for painting in Q4.",
    "Valves and Instrumentation: Valve A-403 (HP Steam line) showed signs of micro-fissures on the upstream flange during dye-penetrant testing. Acoustic emission sensor indicated a faint internal bypass leak when fully closed.",
    "Action Items: Schedule immediate replacement of Valve A-403 as per SOP-MRPL-402. Order 2 units of A-403 valves from vendor (Larsen & Toubro)."
]

for i, bullet_point in enumerate(bullet_points, start=1):
    run = body_shape.text_frame.add_paragraph()
    run.text = bullet_point
    run.font.size = Pt(16)

# Save the presentation
prs.save('Q3_Briefing.pptx')