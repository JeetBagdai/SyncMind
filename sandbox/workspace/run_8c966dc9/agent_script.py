from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Q3 2026 Distillation Unit Inspection Report"
subtitle.text = "Detailed Findings"

# Summary Slide
summary_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(summary_slide_layout)
title = slide.shapes.title
body_shape = slide.placeholders[1]

title.text = "Summary of Detailed Findings"

# Add bullet points to the body
bullets = body_shape.text_frame.add_paragraph()
bullets.text = "Heat Exchangers (HX-101 to HX-104): No significant fouling detected."
bullets.level = 0

bullets = body_shape.text_frame.add_paragraph()
bullets.text = "Piping Network (Line L-205): Minor surface oxidation observed on the external insulation cladding. Recommended for painting in Q4."
bullets.level = 0

bullets = body_shape.text_frame.add_paragraph()
bullets.text = "Valves and Instrumentation: Valve A-403 (HP Steam line) showed signs of micro-fissures on the upstream flange during dye-penetrant testing."
bullets.level = 0

bullets = body_shape.text_frame.add_paragraph()
bullets.text = "Acoustic emission sensor indicated a faint internal bypass leak when fully closed."
bullets.level = 0

# Save the presentation
prs.save('Q3_Briefing.pptx')